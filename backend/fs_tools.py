import os
import re
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from fpdf import FPDF

HEADER_RE = re.compile(r"^[A-Z][A-Z \-/&]{2,}$")

def read_file(filepath):
    try:
        if os.path.isdir(filepath):
            results = []
            for entry in list_files(filepath):
                sub_result = read_file(os.path.join(filepath, entry["name"]))
                if sub_result["success"]:
                    results.append({"file": entry["name"], "content": sub_result["content"]})
            return {"success": True, "files": results}

        ext = os.path.splitext(filepath)[1].lower()

        if ext == ".txt":
            with open(filepath, "r", encoding="utf-8") as f:
                content = f.read()
        elif ext == ".pdf":
            with open(filepath, "rb") as f:
                reader = PdfReader(f)
                content = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext == ".docx":
            doc = Document(filepath)
            content = "\n".join(p.text for p in doc.paragraphs)
        elif ext == ".pptx":
            prs = Presentation(filepath)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in para.runs)
                            if text:
                                lines.append(text)
            content = "\n".join(lines)
        else:
            return {"success": False, "error": f"Unsupported extension: {ext}"}

        stat = os.stat(filepath)
        return {
            "success": True,
            "content": content,
            "metadata": {
                "filename": os.path.basename(filepath),
                "extension": ext,
                "size_bytes": stat.st_size,
                "num_chars": len(content),
            },
        }
    except Exception as e:
        return {"success": False, "error": str(e)}

def list_files(directory, extension=None):
    try:
        files = []
        for name in os.listdir(directory):
            path = os.path.join(directory, name)
            if not os.path.isfile(path):
                continue
            if extension and not name.lower().endswith(extension.lower()):
                continue
            stat = os.stat(path)
            files.append({
                "name": name,
                "size_bytes": stat.st_size,
                "modified": stat.st_mtime,
            })
        return files
    except Exception as e:
        return [{"success": False, "error": str(e)}]

def _split_sections(content):
    """Splits text into (header, [lines]) sections using ALL-CAPS header lines."""
    sections = []
    header = None
    lines = []
    for line in content.splitlines():
        stripped = line.strip()
        if stripped and HEADER_RE.match(stripped) and len(stripped) < 40:
            if header is not None or lines:
                sections.append((header, lines))
            header = stripped
            lines = []
        else:
            lines.append(line)
    sections.append((header, lines))
    return sections


def _write_docx(filepath, content):
    doc = Document()
    for header, lines in _split_sections(content):
        if header:
            doc.add_heading(header, level=2)
        for line in lines:
            if line.strip():
                doc.add_paragraph(line.strip())
    doc.save(filepath)


def _write_pdf(filepath, content):
    pdf = FPDF()
    pdf.set_left_margin(15)
    pdf.set_right_margin(15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)
    for line in content.splitlines():
        safe_line = line.encode("latin-1", "replace").decode("latin-1").strip()
        pdf.write(6, safe_line if safe_line else " ")
        pdf.ln(6)
    pdf.output(filepath)


def _write_pptx(filepath, content):
    prs = Presentation()
    sections = _split_sections(content)

    title_layout = prs.slide_layouts[0]
    first_header, first_lines = sections[0]
    slide = prs.slides.add_slide(title_layout)
    title_text = first_header or (first_lines[0].strip() if first_lines else "Resume")
    slide.shapes.title.text = title_text
    if slide.placeholders[1].has_text_frame:
        slide.placeholders[1].text_frame.text = "\n".join(l.strip() for l in first_lines if l.strip())

    body_layout = prs.slide_layouts[1]
    for header, lines in sections[1:]:
        slide = prs.slides.add_slide(body_layout)
        slide.shapes.title.text = header or "Details"
        body = slide.placeholders[1].text_frame
        text_lines = [l.strip() for l in lines if l.strip()]
        if text_lines:
            body.text = text_lines[0]
            for extra in text_lines[1:]:
                p = body.add_paragraph()
                p.text = extra

    prs.save(filepath)


def write_file(filepath, content):
    try:
        os.makedirs(os.path.dirname(filepath) or ".", exist_ok=True)
        ext = os.path.splitext(filepath)[1].lower()

        if ext == ".txt":
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)
        elif ext == ".docx":
            _write_docx(filepath, content)
        elif ext == ".pdf":
            _write_pdf(filepath, content)
        elif ext == ".pptx":
            _write_pptx(filepath, content)
        else:
            return {"success": False, "error": f"Unsupported extension: {ext}"}

        return {"success": True, "filepath": filepath}
    except Exception as e:
        return {"success": False, "error": str(e)}

def search_in_file(filepath, keyword):
    try:
        if os.path.isdir(filepath):
            all_matches = []
            for entry in list_files(filepath):
                sub_result = search_in_file(os.path.join(filepath, entry["name"]), keyword)
                if sub_result["success"] and sub_result["matches"]:
                    all_matches.append({"file": entry["name"], "matches": sub_result["matches"]})
            return {"success": True, "matches": all_matches}

        result = read_file(filepath)
        if not result["success"]:
            return result

        matches = []
        for i, line in enumerate(result["content"].splitlines(), start=1):
            if keyword.lower() in line.lower():
                matches.append({"line_number": i, "context": line.strip()})

        return {"success": True, "matches": matches}
    except Exception as e:
        return {"success": False, "error": str(e)}
